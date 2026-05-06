import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors (Matching the Dashboard)
BG_COLOR = RGBColor(15, 23, 42)        # #0f172a Dark Blue
TEXT_COLOR = RGBColor(248, 250, 252)   # #f8fafc White
ACCENT_BLUE = RGBColor(96, 165, 250)   # #60a5fa Light Blue
ACCENT_PURPLE = RGBColor(192, 132, 252)# #c084fc Light Purple
MUTED_TEXT = RGBColor(148, 163, 184)   # #94a3b8 Gray

def apply_theme(slide, title_shape, body_shape=None, is_title_slide=False):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    if not is_title_slide:
        left = top = Inches(0)
        width = Inches(10)
        height = Inches(0.15)
        txBox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = ACCENT_PURPLE
        txBox.line.fill.background()
    else:
        txBox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(10), Inches(0.5))
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = ACCENT_BLUE
        txBox.line.fill.background()

    if title_shape and title_shape.has_text_frame:
        for p in title_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if is_title_slide else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = 'Arial'
                run.font.size = Pt(44) if is_title_slide else Pt(36)
                run.font.bold = True
                run.font.color.rgb = ACCENT_BLUE if is_title_slide else ACCENT_PURPLE
                
    if body_shape and body_shape.has_text_frame:
        for p in body_shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if is_title_slide else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = 'Arial'
                run.font.color.rgb = TEXT_COLOR
                if is_title_slide:
                    run.font.size = Pt(20)
                    run.font.color.rgb = MUTED_TEXT

def create_presentation():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Powered Log Analysis"
    subtitle.text = "Automated CI/CD Intelligence with GitHub Actions\nProof of Concept"
    apply_theme(slide, title, subtitle, is_title_slide=True)
    
    # Slide 2: The Goal
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Objective"
    tf = body_shape.text_frame
    tf.text = "Automate the root cause analysis of build failures:"
    for p_text in ["Reduce manual log inspection time.",
                   "Provide actionable fixes immediately on PRs.",
                   "Centralize failure tracking via a dashboard."]:
        p = tf.add_paragraph()
        p.text = p_text
        p.level = 1
    apply_theme(slide, title_shape, body_shape)

    # Slide 3: Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "How It Works"
    tf = body_shape.text_frame
    tf.text = "End-to-End Automation Pipeline:"
    for p_text in ["GitHub Actions triggers on push/PR.",
                   "Failing logs (Build/Test) are captured.",
                   "AI Agent (Gemini/Claude/OpenAI) analyzes context.",
                   "Insights are posted to PR and local Dashboard."]:
        p = tf.add_paragraph()
        p.text = p_text
        p.level = 1
    apply_theme(slide, title_shape, body_shape)

    prs.save('AI_Log_Analysis_POC.pptx')
    print("Presentation saved as AI_Log_Analysis_POC.pptx")

if __name__ == '__main__':
    create_presentation()
